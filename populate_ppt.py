#!/usr/bin/env python3
"""
Script to populate a PowerPoint template with data.
Populates named placeholders: slide_title, role_name, talent_name, and risk_action_table.
"""

from pptx import Presentation
import sys
import os


def find_shape_by_name(slide, name):
    """Find a shape by its name in a slide."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def populate_text_placeholder(shape, text):
    """Populate a text placeholder with the given text."""
    if shape is None:
        return False

    if hasattr(shape, "text_frame"):
        shape.text_frame.text = text
        return True
    elif hasattr(shape, "text"):
        shape.text = text
        return True

    return False


def populate_table(table_shape, data, skip_header=True):
    """
    Populate a table with data.

    Args:
        table_shape: The shape containing the table
        data: List of lists, where each inner list is a row
        skip_header: If True, start populating from row 1 (preserves header row)
    """
    if not table_shape.has_table:
        print(f"Warning: Shape '{table_shape.name}' is not a table")
        return False

    table = table_shape.table
    start_row = 1 if skip_header else 0

    # Check if we have enough rows
    available_rows = len(table.rows) - start_row
    if len(data) > available_rows:
        print(f"Warning: Data has {len(data)} rows but table only has {available_rows} available rows")

    # Populate the table
    for data_idx, row_data in enumerate(data):
        table_row_idx = data_idx + start_row
        if table_row_idx >= len(table.rows):
            break

        for col_idx, cell_value in enumerate(row_data):
            if col_idx >= len(table.columns):
                break

            cell = table.cell(table_row_idx, col_idx)
            cell.text = str(cell_value)

    return True


def populate_presentation(template_path, output_path, data):
    """
    Populate the PowerPoint template with data.

    Args:
        template_path: Path to the template PPTX file
        output_path: Path to save the populated PPTX file
        data: Dictionary containing the data to populate
              Expected keys: 'slide_title', 'role_name', 'talent_name', 'risk_action_table'
    """
    try:
        # Load the presentation
        prs = Presentation(template_path)
        print(f"Loaded template: {template_path}")
        print(f"Total slides: {len(prs.slides)}")

        if len(prs.slides) == 0:
            print("Error: No slides found in the template")
            return False

        slide = prs.slides[0]

        # Populate text placeholders
        placeholders = ['slide_title', 'role_name', 'talent_name']
        for placeholder_name in placeholders:
            if placeholder_name in data:
                shape = find_shape_by_name(slide, placeholder_name)
                if shape:
                    success = populate_text_placeholder(shape, data[placeholder_name])
                    if success:
                        print(f"✓ Populated '{placeholder_name}' with: {data[placeholder_name]}")
                    else:
                        print(f"✗ Failed to populate '{placeholder_name}'")
                else:
                    print(f"✗ Shape '{placeholder_name}' not found in slide")

        # Populate table
        if 'risk_action_table' in data:
            table_shape = find_shape_by_name(slide, 'risk_action_table')
            if table_shape:
                success = populate_table(table_shape, data['risk_action_table'], skip_header=True)
                if success:
                    print(f"✓ Populated table 'risk_action_table'")
                else:
                    print(f"✗ Failed to populate table 'risk_action_table'")
            else:
                print(f"✗ Table shape 'risk_action_table' not found in slide")

        # Save the presentation
        prs.save(output_path)
        print(f"\n✓ Saved populated presentation to: {output_path}")
        return True

    except FileNotFoundError:
        print(f"Error: Template file '{template_path}' not found")
        return False
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """Main function with example usage."""
    template_path = "valueactionplan.pptx"
    output_path = "valueactionplan_populated.pptx"

    if not os.path.exists(template_path):
        print(f"Error: Template file '{template_path}' not found")
        sys.exit(1)

    # Example data - matches CEO Works Value Action Plan format
    data = {
        'slide_title': 'SAMPLE VALUE ACTION PLAN TO MITIGATE RISKS',
        'role_name': 'Chief Marketing Officer',
        'talent_name': 'John Smith',
        'risk_action_table': [
            ['1', 'Alignment of marketing with businesses', 'Align with BUs on the "must-wins" he will prioritize and partner with them to support'],
            ['2', 'Role clarity on swim lanes NA - Global', 'Strengthen partnership with Global Marketing Excellence to drive a fit-for-purpose competency model for NA Zone'],
            ['3', 'Address Capability Gaps in marketing', 'Based on aligned priorities and path; establish the OKRs for his org that indicate value generation and growth signals'],
            ['4', 'Personal networks and ability to influence the organization', 'Finalize executive coaching assignment\nBuild personal relationships across ELT\nPair with identified mentor'],
        ]
    }

    print("Starting population process...\n")
    success = populate_presentation(template_path, output_path, data)

    if success:
        print("\n" + "="*60)
        print("SUCCESS: Presentation populated successfully!")
        print("="*60)
    else:
        print("\n" + "="*60)
        print("FAILED: Could not populate presentation")
        print("="*60)
        sys.exit(1)


if __name__ == "__main__":
    main()