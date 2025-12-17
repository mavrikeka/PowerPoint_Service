#!/usr/bin/env python3
"""
Debug script to print all shape names in a PowerPoint presentation.
This helps identify the names of placeholders and shapes that can be populated.
"""

from pptx import Presentation
import sys


def print_shape_info(prs):
    """Print detailed information about all shapes in the presentation."""
    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"\n{'='*60}")
        print(f"SLIDE {slide_idx}")
        print(f"{'='*60}")

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            print(f"\nShape {shape_idx}:")
            print(f"  Name: {shape.name}")
            print(f"  Type: {shape.shape_type}")

            # Check if it has text
            if hasattr(shape, "text"):
                print(f"  Text: {shape.text[:50]}..." if len(shape.text) > 50 else f"  Text: {shape.text}")

            # Check if it's a placeholder
            if hasattr(shape, "is_placeholder"):
                print(f"  Is Placeholder: {shape.is_placeholder}")
                if shape.is_placeholder:
                    print(f"  Placeholder Type: {shape.placeholder_format.type}")

            # Check if it's a table
            if shape.has_table:
                table = shape.table
                print(f"  Is Table: Yes")
                print(f"  Rows: {len(table.rows)}, Columns: {len(table.columns)}")

            # Check if it has text frame
            if hasattr(shape, "text_frame"):
                print(f"  Has Text Frame: Yes")


def main():
    """Main function to run the debug script."""
    if len(sys.argv) < 2:
        print("Usage: python debug_shapes.py <path_to_pptx>")
        print("Example: python debug_shapes.py valueactionplan.pptx")
        sys.exit(1)

    pptx_file = sys.argv[1]

    try:
        print(f"Opening presentation: {pptx_file}")
        prs = Presentation(pptx_file)
        print(f"Total slides: {len(prs.slides)}")
        print_shape_info(prs)

    except FileNotFoundError:
        print(f"Error: File '{pptx_file}' not found.")
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
