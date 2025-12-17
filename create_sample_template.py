#!/usr/bin/env python3
"""
Script to create a sample PowerPoint template with named placeholders.
This creates a template that can be populated using populate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_sample_template(output_path='valueactionplan_template.pptx'):
    """Create a sample PowerPoint template with named shapes."""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add slide title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_box.name = "slide_title"
    title_frame = title_box.text_frame
    title_frame.text = "[SLIDE TITLE]"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER

    # Add role name label and placeholder
    role_label = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(2), Inches(0.5)
    )
    role_label_frame = role_label.text_frame
    role_label_frame.text = "Role:"
    role_label_frame.paragraphs[0].font.size = Pt(14)
    role_label_frame.paragraphs[0].font.bold = True

    role_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(1.8), Inches(7), Inches(0.5)
    )
    role_box.name = "role_name"
    role_frame = role_box.text_frame
    role_frame.text = "[ROLE NAME]"
    role_frame.paragraphs[0].font.size = Pt(14)

    # Add talent name label and placeholder
    talent_label = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(2), Inches(0.5)
    )
    talent_label_frame = talent_label.text_frame
    talent_label_frame.text = "Talent:"
    talent_label_frame.paragraphs[0].font.size = Pt(14)
    talent_label_frame.paragraphs[0].font.bold = True

    talent_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(2.5), Inches(7), Inches(0.5)
    )
    talent_box.name = "talent_name"
    talent_frame = talent_box.text_frame
    talent_frame.text = "[TALENT NAME]"
    talent_frame.paragraphs[0].font.size = Pt(14)

    # Add table for risk/action items
    rows, cols = 5, 4  # 1 header row + 4 data rows, 4 columns
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(3.5),
        Inches(9), Inches(3.5)
    )
    table_shape.name = "risk_action_table"
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.5)  # Risk
    table.columns[1].width = Inches(2.5)  # Action
    table.columns[2].width = Inches(2)    # Owner
    table.columns[3].width = Inches(2)    # Due Date

    # Format header row
    header_texts = ['Risk', 'Action', 'Owner', 'Due Date']
    for col_idx, header_text in enumerate(header_texts):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Add placeholder text to data rows
    for row_idx in range(1, rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = f"[{header_texts[col_idx]} {row_idx}]"
            cell.text_frame.paragraphs[0].font.size = Pt(11)

    # Save the presentation
    prs.save(output_path)
    print(f"✓ Sample template created: {output_path}")
    print("\nThis template includes:")
    print("  - slide_title: Main title placeholder")
    print("  - role_name: Role/position placeholder")
    print("  - talent_name: Person name placeholder")
    print("  - risk_action_table: 4-column table (Risk, Action, Owner, Due Date)")
    print("\nYou can now:")
    print("  1. Open this file in PowerPoint to customize the design")
    print("  2. Use populate_ppt.py to fill it with data")
    print("  3. Or replace valueactionplan.pptx with this file")


if __name__ == "__main__":
    create_sample_template()
