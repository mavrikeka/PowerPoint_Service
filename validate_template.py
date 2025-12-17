#!/usr/bin/env python3
"""
Template Validator Utility

Inspects PowerPoint templates to:
- Check auto-fit settings on text shapes
- Generate JSON definition for population
- Validate template structure
- Auto-fix common issues
- Generate HTML reports

Usage:
  python validate_template.py template.pptx
  python validate_template.py template.pptx --fix
  python validate_template.py template.pptx --export-html
"""

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
import json
import sys
from pathlib import Path
from datetime import datetime


class TemplateValidator:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)
        self.issues = []
        self.warnings = []
        self.slides_info = []

    def inspect_template(self):
        """Inspect all slides in the template."""
        for slide_idx, slide in enumerate(self.prs.slides):
            slide_info = {
                'slide_index': slide_idx,
                'text_fields': [],
                'tables': [],
                'other_shapes': []
            }

            for shape in slide.shapes:
                shape_data = self._inspect_shape(shape, slide_idx)

                if shape_data['type'] == 'text':
                    slide_info['text_fields'].append(shape_data)
                elif shape_data['type'] == 'table':
                    slide_info['tables'].append(shape_data)
                else:
                    slide_info['other_shapes'].append(shape_data)

            self.slides_info.append(slide_info)

    def _inspect_shape(self, shape, slide_idx):
        """Inspect a single shape and return its properties."""
        shape_data = {
            'name': shape.name,
            'type': 'unknown'
        }

        # Check if it's a table
        if shape.has_table:
            table = shape.table
            shape_data['type'] = 'table'
            shape_data['rows'] = len(table.rows)
            shape_data['columns'] = len(table.columns)
            shape_data['headers'] = [cell.text for cell in table.rows[0].cells] if len(table.rows) > 0 else []
            return shape_data

        # Check if it has text frame
        if hasattr(shape, 'text_frame'):
            shape_data['type'] = 'text'
            text_frame = shape.text_frame

            # Check auto-fit setting
            autofit_type = text_frame.auto_size
            shape_data['autofit'] = self._get_autofit_status(autofit_type)

            # Check if auto-fit is disabled
            if autofit_type == MSO_AUTO_SIZE.NONE:
                self.warnings.append({
                    'slide': slide_idx,
                    'shape': shape.name,
                    'issue': 'Auto-fit is disabled',
                    'recommendation': 'Enable auto-fit to prevent text overflow'
                })

            # Get font info
            if text_frame.paragraphs:
                first_para = text_frame.paragraphs[0]
                if first_para.runs:
                    font = first_para.runs[0].font
                    shape_data['font_size'] = font.size.pt if font.size else 'Unknown'
                else:
                    shape_data['font_size'] = 'Unknown'
            else:
                shape_data['font_size'] = 'Unknown'

            shape_data['text_preview'] = shape.text[:50] if shape.text else '[Empty]'

            return shape_data

        return shape_data

    def _get_autofit_status(self, autofit_type):
        """Convert autofit type to readable string."""
        if autofit_type == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
            return 'SHAPE_TO_FIT_TEXT'
        elif autofit_type == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
            return 'TEXT_TO_FIT_SHAPE'
        elif autofit_type == MSO_AUTO_SIZE.NONE:
            return 'OFF'
        else:
            return 'UNKNOWN'

    def generate_json_definition(self):
        """Generate example JSON for populating this template."""
        definition = {
            'slides': []
        }

        for slide_info in self.slides_info:
            slide_data = {
                'slide_index': slide_info['slide_index'],
                'data': {}
            }

            # Add text fields
            for field in slide_info['text_fields']:
                slide_data['data'][field['name']] = f"[{field['name']}]"

            # Add tables
            for table in slide_info['tables']:
                example_row = ["Col " + str(i+1) for i in range(table['columns'])]
                slide_data['data'][table['name']] = [example_row] * 3

            definition['slides'].append(slide_data)

        return definition

    def print_report(self):
        """Print a comprehensive validation report."""
        print("=" * 65)
        print(f"TEMPLATE VALIDATION REPORT: {Path(self.pptx_path).name}")
        print("=" * 65)
        print()
        print(f"📊 SLIDES: {len(self.prs.slides)} slide(s) found")
        print()

        # Print each slide
        for slide_info in self.slides_info:
            print("━" * 65)
            print(f"SLIDE {slide_info['slide_index']}")
            print("━" * 65)
            print()

            # Text fields
            if slide_info['text_fields']:
                print("TEXT FIELDS:")
                for field in slide_info['text_fields']:
                    autofit_symbol = "✓" if field['autofit'] != 'OFF' else "✗"
                    print(f"  {autofit_symbol} {field['name']:<20} [AutoFit: {field['autofit']:<20}] Font: {field['font_size']}")
                print()

            # Tables
            if slide_info['tables']:
                print("TABLES:")
                for table in slide_info['tables']:
                    print(f"  ✓ {table['name']:<20} [{table['columns']} columns × {table['rows']} rows]")
                    if table['headers']:
                        print(f"    Headers: {', '.join(table['headers'])}")
                print()

            # Other shapes
            if slide_info['other_shapes']:
                print("OTHER SHAPES:")
                for shape in slide_info['other_shapes']:
                    print(f"  • {shape['name']} ({shape['type']})")
                print()

        # Print warnings
        print("━" * 65)
        if self.warnings:
            print(f"ISSUES FOUND: {len(self.warnings)}")
            print("━" * 65)
            print()
            for warning in self.warnings:
                print(f"⚠️  Slide {warning['slide']}: '{warning['shape']}' - {warning['issue']}")
                print(f"    💡 {warning['recommendation']}")
                print()
        else:
            print("NO ISSUES FOUND ✓")
            print("━" * 65)
            print()

        # Print JSON definition
        print("━" * 65)
        print("JSON DEFINITION")
        print("━" * 65)
        print()
        print("Example JSON for this template:")
        print()
        definition = self.generate_json_definition()
        print(json.dumps(definition, indent=2))
        print()

        # Save JSON
        json_path = Path(self.pptx_path).stem + '_definition.json'
        with open(json_path, 'w') as f:
            json.dump(definition, f, indent=2)
        print(f"✓ JSON saved to: {json_path}")
        print()

        # Print options
        if self.warnings:
            print("━" * 65)
            print("OPTIONS")
            print("━" * 65)
            print()
            print("Run with --fix to automatically enable auto-fit on all text fields:")
            print(f"  python validate_template.py {Path(self.pptx_path).name} --fix")
            print()

    def auto_fix(self, output_path=None):
        """Auto-fix common issues and save to a new file."""
        if output_path is None:
            output_path = Path(self.pptx_path).stem + '_fixed.pptx'

        fixed_count = 0

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    text_frame = shape.text_frame

                    # Enable auto-fit if it's disabled
                    if text_frame.auto_size == MSO_AUTO_SIZE.NONE:
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        fixed_count += 1

        self.prs.save(output_path)

        print("=" * 65)
        print("AUTO-FIX COMPLETE")
        print("=" * 65)
        print()
        print(f"✓ Fixed {fixed_count} text field(s)")
        print(f"✓ Saved to: {output_path}")
        print()
        print("All text fields now have auto-fit enabled!")
        print()

        return fixed_count

    def generate_html_report(self, output_path=None):
        """Generate an HTML report."""
        if output_path is None:
            output_path = Path(self.pptx_path).stem + '_report.html'

        html = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Template Validation Report - {Path(self.pptx_path).name}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Arial, sans-serif;
            max-width: 1200px;
            margin: 0 auto;
            padding: 40px 20px;
            background: #f5f5f5;
        }}
        .header {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 30px;
            border-radius: 12px;
            margin-bottom: 30px;
        }}
        .header h1 {{ margin: 0 0 10px 0; }}
        .header p {{ margin: 0; opacity: 0.9; }}
        .stats {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: 20px;
            margin-bottom: 30px;
        }}
        .stat-card {{
            background: white;
            padding: 20px;
            border-radius: 8px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        }}
        .stat-card h3 {{ margin: 0 0 10px 0; color: #667eea; font-size: 14px; }}
        .stat-card .value {{ font-size: 32px; font-weight: bold; color: #333; }}
        .slide {{
            background: white;
            padding: 30px;
            border-radius: 12px;
            margin-bottom: 20px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        }}
        .slide h2 {{
            margin: 0 0 20px 0;
            padding-bottom: 10px;
            border-bottom: 2px solid #667eea;
        }}
        .section {{ margin-bottom: 20px; }}
        .section h3 {{
            color: #667eea;
            margin: 0 0 10px 0;
            font-size: 16px;
        }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin-top: 10px;
        }}
        th, td {{
            text-align: left;
            padding: 12px;
            border-bottom: 1px solid #eee;
        }}
        th {{
            background: #f8f9fa;
            font-weight: 600;
            color: #333;
        }}
        .status-ok {{ color: #28a745; }}
        .status-warning {{ color: #ffc107; }}
        .badge {{
            display: inline-block;
            padding: 4px 8px;
            border-radius: 4px;
            font-size: 12px;
            font-weight: 600;
        }}
        .badge-success {{ background: #d4edda; color: #155724; }}
        .badge-warning {{ background: #fff3cd; color: #856404; }}
        .warnings {{
            background: #fff3cd;
            border-left: 4px solid #ffc107;
            padding: 20px;
            border-radius: 8px;
            margin-bottom: 30px;
        }}
        .warnings h3 {{
            color: #856404;
            margin-top: 0;
        }}
        .warning-item {{
            margin: 10px 0;
            padding: 10px;
            background: white;
            border-radius: 4px;
        }}
        .json-preview {{
            background: #f8f9fa;
            border: 1px solid #dee2e6;
            border-radius: 8px;
            padding: 20px;
            margin-top: 20px;
        }}
        .json-preview h3 {{ margin-top: 0; }}
        pre {{
            background: white;
            padding: 15px;
            border-radius: 4px;
            overflow-x: auto;
        }}
        .footer {{
            text-align: center;
            color: #666;
            margin-top: 40px;
            padding-top: 20px;
            border-top: 1px solid #dee2e6;
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>📊 Template Validation Report</h1>
        <p>{Path(self.pptx_path).name}</p>
        <p style="font-size: 14px; margin-top: 10px;">Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
    </div>

    <div class="stats">
        <div class="stat-card">
            <h3>Total Slides</h3>
            <div class="value">{len(self.prs.slides)}</div>
        </div>
        <div class="stat-card">
            <h3>Text Fields</h3>
            <div class="value">{sum(len(s['text_fields']) for s in self.slides_info)}</div>
        </div>
        <div class="stat-card">
            <h3>Tables</h3>
            <div class="value">{sum(len(s['tables']) for s in self.slides_info)}</div>
        </div>
        <div class="stat-card">
            <h3>Issues Found</h3>
            <div class="value" style="color: {'#ffc107' if self.warnings else '#28a745'}">{len(self.warnings)}</div>
        </div>
    </div>
"""

        # Warnings section
        if self.warnings:
            html += """
    <div class="warnings">
        <h3>⚠️ Issues Found</h3>
"""
            for warning in self.warnings:
                html += f"""
        <div class="warning-item">
            <strong>Slide {warning['slide']}: {warning['shape']}</strong><br>
            {warning['issue']}<br>
            <em>💡 {warning['recommendation']}</em>
        </div>
"""
            html += """
    </div>
"""

        # Slides
        for slide_info in self.slides_info:
            html += f"""
    <div class="slide">
        <h2>Slide {slide_info['slide_index']}</h2>
"""

            # Text fields
            if slide_info['text_fields']:
                html += """
        <div class="section">
            <h3>Text Fields</h3>
            <table>
                <tr>
                    <th>Shape Name</th>
                    <th>Auto-Fit</th>
                    <th>Font Size</th>
                    <th>Preview</th>
                </tr>
"""
                for field in slide_info['text_fields']:
                    status_class = 'status-ok' if field['autofit'] != 'OFF' else 'status-warning'
                    badge_class = 'badge-success' if field['autofit'] != 'OFF' else 'badge-warning'
                    html += f"""
                <tr>
                    <td>{field['name']}</td>
                    <td><span class="badge {badge_class}">{field['autofit']}</span></td>
                    <td>{field['font_size']}</td>
                    <td style="font-size: 12px; color: #666;">{field['text_preview']}</td>
                </tr>
"""
                html += """
            </table>
        </div>
"""

            # Tables
            if slide_info['tables']:
                html += """
        <div class="section">
            <h3>Tables</h3>
            <table>
                <tr>
                    <th>Table Name</th>
                    <th>Dimensions</th>
                    <th>Headers</th>
                </tr>
"""
                for table in slide_info['tables']:
                    headers = ', '.join(table['headers']) if table['headers'] else 'None'
                    html += f"""
                <tr>
                    <td>{table['name']}</td>
                    <td>{table['columns']} × {table['rows']}</td>
                    <td style="font-size: 12px;">{headers}</td>
                </tr>
"""
                html += """
            </table>
        </div>
"""

            html += """
    </div>
"""

        # JSON definition
        definition = self.generate_json_definition()
        html += f"""
    <div class="slide">
        <h2>📝 JSON Definition</h2>
        <p>Use this JSON structure to populate the template:</p>
        <pre>{json.dumps(definition, indent=2)}</pre>
    </div>

    <div class="footer">
        <p>Generated by PowerPoint Template Validator</p>
        <p style="font-size: 12px; margin-top: 10px;">
            To fix issues automatically: <code>python validate_template.py {Path(self.pptx_path).name} --fix</code>
        </p>
    </div>
</body>
</html>
"""

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)

        print("=" * 65)
        print("HTML REPORT GENERATED")
        print("=" * 65)
        print()
        print(f"✓ Report saved to: {output_path}")
        print(f"  Open in browser: file://{Path(output_path).absolute()}")
        print()


def main():
    if len(sys.argv) < 2:
        print("Usage: python validate_template.py <template.pptx> [options]")
        print()
        print("Options:")
        print("  --fix           Auto-fix issues and save to new file")
        print("  --export-html   Generate HTML report")
        print()
        print("Examples:")
        print("  python validate_template.py template.pptx")
        print("  python validate_template.py template.pptx --fix")
        print("  python validate_template.py template.pptx --export-html")
        print("  python validate_template.py template.pptx --fix --export-html")
        sys.exit(1)

    pptx_path = sys.argv[1]

    if not Path(pptx_path).exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    # Create validator
    validator = TemplateValidator(pptx_path)

    # Inspect template
    validator.inspect_template()

    # Print report
    validator.print_report()

    # Auto-fix if requested
    if '--fix' in sys.argv:
        validator.auto_fix()

    # Generate HTML report if requested
    if '--export-html' in sys.argv:
        validator.generate_html_report()


if __name__ == "__main__":
    main()
