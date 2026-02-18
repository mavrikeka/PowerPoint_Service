#!/usr/bin/env python3
"""
FastAPI service for populating PowerPoint templates.
Upload a template + data, receive a populated PowerPoint file.
"""

from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import Response, FileResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
import json
import tempfile
import os
from pathlib import Path
from typing import Optional
import shutil
from copy import deepcopy

app = FastAPI(
    title="PowerPoint Population Service",
    description="Upload a PowerPoint template and data to generate a populated presentation",
    version="1.0.0"
)

# Add CORS middleware to allow requests from your application
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Configure this for production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def find_shape_by_name(slide, name):
    """
    Find a shape by its name in a slide.

    Handles duplicate shape names with _N suffix (e.g., "object 39_1" for the second "object 39").
    """
    # Check if the name has a _N suffix for duplicate handling
    import re
    match = re.match(r'^(.+)_(\d+)$', name)

    if match:
        # Extract base name and index
        base_name = match.group(1)
        index = int(match.group(2))

        # Find all shapes with the base name
        matching_shapes = [s for s in slide.shapes if s.name == base_name]

        # Return the shape at the specified index (if it exists)
        if index < len(matching_shapes):
            return matching_shapes[index]

    # Direct name match (no suffix)
    for shape in slide.shapes:
        if shape.name == name:
            return shape

    return None


def _capture_font_props(font):
    """Capture font properties from a python-pptx font object."""
    return {
        'name': font.name,
        'size': font.size,
        'bold': font.bold,
        'italic': font.italic,
        'underline': font.underline,
        'color_type': font.color.type if hasattr(font.color, 'type') else None,
        'color_rgb': font.color.rgb if (hasattr(font.color, 'type') and font.color.type == 1) else None,
        'color_theme': font.color.theme_color if (hasattr(font.color, 'theme_color')) else None,
    }


def _apply_font_props(font, props):
    """Apply saved font properties to a python-pptx font object."""
    if props['name'] is not None:
        font.name = props['name']
    if props['size'] is not None:
        font.size = props['size']
    if props['bold'] is not None:
        font.bold = props['bold']
    if props['italic'] is not None:
        font.italic = props['italic']
    if props['underline'] is not None:
        font.underline = props['underline']
    if props['color_rgb'] is not None:
        font.color.rgb = props['color_rgb']
    elif props['color_theme'] is not None:
        font.color.theme_color = props['color_theme']


def populate_text_placeholder(shape, text):
    """
    Populate a text placeholder with the given text while preserving formatting.

    Captures font properties from the template's first paragraph and applies them
    to ALL paragraphs created after population (text with newlines creates multiple
    paragraphs). Without this, only paragraph 0 gets the correct font — subsequent
    paragraphs fall back to the PowerPoint theme default (typically 18pt).
    """
    if shape is None:
        return False

    if hasattr(shape, "text_frame"):
        text_frame = shape.text_frame

        # Check if there's existing content with formatting to preserve
        if len(text_frame.paragraphs) > 0:
            first_para = text_frame.paragraphs[0]

            # Capture font from the first run if it exists, otherwise paragraph-level
            saved_font_props = None
            if len(first_para.runs) > 0:
                saved_font_props = _capture_font_props(first_para.runs[0].font)
            else:
                # No runs (empty placeholder) - capture paragraph-level formatting
                saved_font_props = _capture_font_props(first_para.font)

            # Replace the text (may create multiple paragraphs if text contains newlines)
            text_frame.text = text

            # Restore formatting to ALL paragraphs, not just the first.
            # When text contains '\n', text_frame.text = text creates one paragraph
            # per line. Only para[0] inherits the restored run properties; every
            # subsequent paragraph gets a bare run with no font attributes, causing
            # PowerPoint to fall back to the theme default (usually 18pt).
            if saved_font_props:
                for new_para in text_frame.paragraphs:
                    if len(new_para.runs) > 0:
                        _apply_font_props(new_para.runs[0].font, saved_font_props)
        else:
            # No existing content, just set the text
            text_frame.text = text

        return True
    elif hasattr(shape, "text"):
        shape.text = text
        return True

    return False


def populate_table(table_shape, data, skip_header=True):
    """
    Populate a table with data while preserving cell formatting.

    Captures font properties from each template cell and applies them to ALL
    paragraphs in the populated cell. Table cells with multi-line content
    (newlines in data) create multiple paragraphs; without restoring props to
    every paragraph, subsequent ones fall back to the PowerPoint theme default.

    Args:
        table_shape: The shape containing the table
        data: List of lists, where each inner list is a row
        skip_header: If True, start populating from row 1 (preserves header row)
    """
    if not table_shape.has_table:
        return False

    table = table_shape.table
    start_row = 1 if skip_header else 0

    # Populate the table with new data
    for data_idx, row_data in enumerate(data):
        table_row_idx = data_idx + start_row
        if table_row_idx >= len(table.rows):
            break

        for col_idx, cell_value in enumerate(row_data):
            if col_idx >= len(table.columns):
                break

            cell = table.cell(table_row_idx, col_idx)

            # Capture font from first run if present, otherwise paragraph-level
            saved_font_props = None
            if cell.text_frame and len(cell.text_frame.paragraphs) > 0:
                first_para = cell.text_frame.paragraphs[0]
                if len(first_para.runs) > 0:
                    saved_font_props = _capture_font_props(first_para.runs[0].font)
                else:
                    saved_font_props = _capture_font_props(first_para.font)

            # Set the cell text (may create multiple paragraphs if value contains newlines)
            cell.text = str(cell_value)

            # Restore formatting to ALL paragraphs in the cell
            if saved_font_props and cell.text_frame:
                for new_para in cell.text_frame.paragraphs:
                    if len(new_para.runs) > 0:
                        _apply_font_props(new_para.runs[0].font, saved_font_props)

    # Clear any remaining rows that weren't overwritten
    rows_populated = len(data) + start_row
    for row_idx in range(rows_populated, len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""

    return True


def populate_single_slide(prs, slide, data: dict):
    """
    Populate a single slide with data.

    Args:
        prs: Presentation object
        slide: Slide object to populate
        data: Dictionary containing the data to populate

    Returns:
        List of populated field names
    """
    populated_fields = []

    # Populate text placeholders
    for key, value in data.items():
        # Skip table data (lists will be handled in the table section)
        if isinstance(value, list):
            continue

        shape = find_shape_by_name(slide, key)
        if shape:
            success = populate_text_placeholder(shape, str(value))
            if success:
                populated_fields.append(key)

    # Populate tables (any field with list data, regardless of naming)
    for key, value in data.items():
        if isinstance(value, list):
            table_shape = find_shape_by_name(slide, key)
            if table_shape:
                # Since extracted data includes the full table with headers, don't skip
                skip_header = data.get(f"{key}_skip_header", False)
                success = populate_table(table_shape, value, skip_header=skip_header)
                if success:
                    populated_fields.append(key)

    return populated_fields


def populate_multi_slide(template_file, output_file, slides_data: list):
    """
    Populate multiple slides from a template.

    Supports intelligent slide duplication:
    1. First occurrence of slide_index: populate that template slide
    2. Second+ occurrences of same slide_index: duplicate that template slide and populate
    3. If slide_index doesn't exist in template: duplicate slide 0

    This allows creating multiple copies from different slide templates.
    Example: 3 slide templates, 3 data items each = 9 total slides

    Args:
        template_file: Path to the template PPTX file
        output_file: Path to save the populated PPTX file
        slides_data: List of dicts with 'slide_index' and 'data' keys

    Returns:
        Tuple of (success: bool, message: str)
    """
    try:
        # Load the template presentation directly
        prs = Presentation(template_file)

        if len(prs.slides) == 0:
            return False, "Template has no slides"

        total_fields = 0
        slides_created = 0

        # Track which slide_index values we've seen (for duplication logic)
        slide_index_usage = {}

        # Process each slide data
        for item in slides_data:
            slide_index = item.get('slide_index', 0)
            data = item.get('data', {})

            # Track how many times we've seen this slide_index
            if slide_index not in slide_index_usage:
                slide_index_usage[slide_index] = 0
            else:
                slide_index_usage[slide_index] += 1

            occurrence_count = slide_index_usage[slide_index]

            # Determine which slide to use or duplicate
            if slide_index < len(prs.slides) and occurrence_count == 0:
                # First occurrence - use the template slide at this index
                slide = prs.slides[slide_index]
            else:
                # Need to duplicate - either it's a repeat or slide doesn't exist
                # Determine which slide to use as the source
                if slide_index < len(prs.slides):
                    # Duplicate the template slide at this index
                    source_slide = prs.slides[slide_index]
                else:
                    # slide_index doesn't exist in template, duplicate slide 0
                    source_slide = prs.slides[0]

                slide_layout = source_slide.slide_layout

                # Create a new slide with the same layout
                new_slide = prs.slides.add_slide(slide_layout)

                # Remove default placeholder shapes added by add_slide
                # These are unwanted "Title 1", "Content Placeholder 2", etc.
                shapes_to_remove = []
                for shape in new_slide.shapes:
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        shapes_to_remove.append(shape)

                for shape in shapes_to_remove:
                    sp = shape.element
                    sp.getparent().remove(sp)

                # Copy all shapes from source slide to new slide
                for shape in source_slide.shapes:
                    # Get the shape element
                    el = shape.element
                    # Create a copy of the shape element
                    newel = deepcopy(el)
                    # Add it to the new slide
                    new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

                slide = new_slide
                slides_created += 1

            # Populate the slide
            populated_fields = populate_single_slide(prs, slide, data)
            total_fields += len(populated_fields)

        # Save the modified presentation
        prs.save(output_file)

        if slides_created > 0:
            message = f"Successfully created {slides_created} new slide(s), populated {len(slides_data)} total slide(s), modified {total_fields} total fields"
        else:
            message = f"Successfully populated {len(slides_data)} slide(s), modified {total_fields} total fields"

        return True, message

    except Exception as e:
        return False, f"Error processing multi-slide template: {str(e)}"


def populate_presentation_from_data(template_file, output_file, data: dict, slide_index: int = 0):
    """
    Populate a PowerPoint template with data.

    Args:
        template_file: Path to the template PPTX file
        output_file: Path to save the populated PPTX file
        data: Dictionary containing the data to populate
        slide_index: Which slide to populate (default: 0 for first slide)

    Returns:
        Tuple of (success: bool, message: str)
    """
    try:
        # Load the presentation
        prs = Presentation(template_file)

        if len(prs.slides) == 0:
            return False, "Template has no slides"

        if slide_index >= len(prs.slides):
            return False, f"Slide index {slide_index} out of range (template has {len(prs.slides)} slides)"

        slide = prs.slides[slide_index]
        populated_fields = []

        # Populate text placeholders
        for key, value in data.items():
            # Skip table data for now
            if key.endswith('_table') or isinstance(value, list):
                continue

            shape = find_shape_by_name(slide, key)
            if shape:
                success = populate_text_placeholder(shape, str(value))
                if success:
                    populated_fields.append(key)

        # Populate tables
        for key, value in data.items():
            if key.endswith('_table') and isinstance(value, list):
                table_shape = find_shape_by_name(slide, key)
                if table_shape:
                    # Check if first row should be treated as header
                    skip_header = data.get(f"{key}_skip_header", True)
                    success = populate_table(table_shape, value, skip_header=skip_header)
                    if success:
                        populated_fields.append(key)

        # Save the presentation
        prs.save(output_file)

        message = f"Successfully populated {len(populated_fields)} fields: {', '.join(populated_fields)}"
        return True, message

    except Exception as e:
        return False, f"Error processing template: {str(e)}"


@app.get("/")
async def root():
    """Serve the test UI."""
    index_path = Path(__file__).parent / "index.html"
    if index_path.exists():
        return FileResponse(index_path, media_type="text/html")
    else:
        return HTMLResponse("""
        <html>
            <body>
                <h1>PowerPoint Population Service</h1>
                <p>Service is running. API endpoint: POST /populate-pptx</p>
                <p>For API documentation, see <a href="https://github.com/CEO-Works/Create_Powerpoint">GitHub</a></p>
            </body>
        </html>
        """)


@app.get("/extract.html")
async def extract_ui():
    """Serve the extraction UI."""
    extract_path = Path(__file__).parent / "extract.html"
    if extract_path.exists():
        return FileResponse(extract_path, media_type="text/html")
    else:
        return HTMLResponse("<h1>Extract UI not found</h1>", status_code=404)


@app.get("/health")
async def health():
    """Health check for Railway."""
    return {"status": "healthy"}


@app.post("/populate-pptx")
async def populate_pptx(
    template: UploadFile = File(..., description="PowerPoint template file (.pptx)"),
    data: str = Form(..., description="JSON string containing field names and values to populate"),
    slide_index: Optional[int] = Form(0, description="Slide index to populate (default: 0, ignored for multi-slide format)"),
    output_filename: Optional[str] = Form("output.pptx", description="Name for the output file")
):
    """
    Populate a PowerPoint template with data.

    Supports two formats:

    1. Single-slide format (backward compatible):
    {
        "slide_title": "My Title",
        "role_name": "Software Engineer",
        "talent_name": "John Doe",
        "risk_action_table": [
            ["Risk 1", "Action 1", "Owner 1", "Date 1"],
            ["Risk 2", "Action 2", "Owner 2", "Date 2"]
        ]
    }

    2. Multi-slide format:
    {
        "slides": [
            {
                "slide_index": 0,
                "data": {
                    "slide_title": "Title 1",
                    "role_name": "Engineer"
                }
            },
            {
                "slide_index": 0,
                "data": {
                    "slide_title": "Title 2",
                    "role_name": "Manager"
                }
            }
        ]
    }
    """
    # Validate template file
    if not template.filename.endswith('.pptx'):
        raise HTTPException(status_code=400, detail="Template must be a .pptx file")

    # Parse JSON data
    try:
        data_dict = json.loads(data)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=400, detail=f"Invalid JSON data: {str(e)}")

    # Create temporary directory for processing
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)

        # Save uploaded template
        template_path = temp_dir_path / "template.pptx"
        with open(template_path, "wb") as f:
            shutil.copyfileobj(template.file, f)

        # Generate output path
        output_path = temp_dir_path / "output.pptx"

        # Detect format and populate accordingly
        if "slides" in data_dict:
            # Multi-slide format (Option 1)
            success, message = populate_multi_slide(
                str(template_path),
                str(output_path),
                data_dict["slides"]
            )
        else:
            # Single-slide format (backward compatible)
            success, message = populate_presentation_from_data(
                str(template_path),
                str(output_path),
                data_dict,
                slide_index
            )

        if not success:
            raise HTTPException(status_code=400, detail=message)

        # Read the file into memory before temp dir is cleaned up
        with open(output_path, "rb") as f:
            file_content = f.read()

        # Return the populated file
        return Response(
            content=file_content,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{output_filename}"',
                "X-Population-Message": message
            }
        )


@app.post("/extract-data")
async def extract_data(
    presentation: UploadFile = File(..., description="PowerPoint file (.pptx) to extract data from"),
    slide_index: Optional[int] = Form(None, description="Slide index to extract from (default: extract all slides)"),
    extract_all: Optional[bool] = Form(False, description="Extract data from all slides")
):
    """
    Extract data from a PowerPoint presentation to JSON format.

    This is the reverse operation of /populate-pptx - it reads a PowerPoint
    file and extracts all text and table data into JSON format.

    Use cases:
    - Convert existing presentations to data
    - Create JSON templates from real presentations
    - Extract data for analysis or migration

    Returns JSON with shape names as keys and extracted content as values.
    """
    # Validate file
    if not presentation.filename.endswith('.pptx'):
        raise HTTPException(status_code=400, detail="File must be a .pptx file")

    # Create temporary directory for processing
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)

        # Save uploaded file
        pptx_path = temp_dir_path / "presentation.pptx"
        with open(pptx_path, "wb") as f:
            shutil.copyfileobj(presentation.file, f)

        try:
            prs = Presentation(pptx_path)

            if len(prs.slides) == 0:
                raise HTTPException(status_code=400, detail="No slides found in presentation")

            # Extract from all slides or specific slide
            if extract_all or slide_index is None:
                # Extract from all slides in populate-ready format
                all_slides_data = []

                for slide_idx, slide in enumerate(prs.slides):
                    slide_shapes = []  # Ordered array of shapes

                    for shape in slide.shapes:
                        shape_data = {
                            "name": shape.name,
                            "type": None,
                            "content": None
                        }

                        if shape.has_table:
                            # Extract table data
                            shape_data["type"] = "table"
                            table = shape.table
                            table_data = []
                            for row in table.rows:
                                row_data = [cell.text for cell in row.cells]
                                table_data.append(row_data)
                            shape_data["content"] = table_data
                            slide_shapes.append(shape_data)
                        else:
                            # Extract text data
                            text = None
                            if hasattr(shape, "text_frame"):
                                text = shape.text_frame.text
                            elif hasattr(shape, "text"):
                                text = shape.text

                            if text:
                                shape_data["type"] = "text"
                                shape_data["content"] = text
                                slide_shapes.append(shape_data)

                    all_slides_data.append({
                        "slide_index": slide_idx,
                        "shapes": slide_shapes
                    })

                # Return populate-ready multi-slide format
                return {
                    "slides": all_slides_data
                }

            else:
                # Extract from specific slide
                if slide_index >= len(prs.slides):
                    raise HTTPException(
                        status_code=400,
                        detail=f"Slide index {slide_index} out of range (presentation has {len(prs.slides)} slides)"
                    )

                slide = prs.slides[slide_index]
                slide_shapes = []  # Ordered array of shapes

                for shape in slide.shapes:
                    shape_data = {
                        "name": shape.name,
                        "type": None,
                        "content": None
                    }

                    if shape.has_table:
                        # Extract table data
                        shape_data["type"] = "table"
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        shape_data["content"] = table_data
                        slide_shapes.append(shape_data)
                    else:
                        # Extract text data
                        text = None
                        if hasattr(shape, "text_frame"):
                            text = shape.text_frame.text
                        elif hasattr(shape, "text"):
                            text = shape.text

                        if text:
                            shape_data["type"] = "text"
                            shape_data["content"] = text
                            slide_shapes.append(shape_data)

                # Return ordered array of shapes
                return slide_shapes

        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error extracting data: {str(e)}")


if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
