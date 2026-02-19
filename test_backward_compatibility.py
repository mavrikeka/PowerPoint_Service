#!/usr/bin/env python3
"""
Test backward compatibility - plain text should work exactly as before.
"""

import sys
import os
from pathlib import Path

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from service import populate_text_placeholder, populate_table
import tempfile


def create_test_presentation():
    """Create a simple test presentation with a text placeholder."""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add a text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(5)
    height = Inches(1)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = "test_placeholder"

    # Set initial formatting
    text_frame = textbox.text_frame
    text_frame.text = "Template Text"

    # Apply specific font properties
    para = text_frame.paragraphs[0]
    run = para.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(14)
    run.font.bold = False
    run.font.italic = False

    # Add a table
    rows, cols = 3, 3
    left = Inches(1)
    top = Inches(3)
    width = Inches(5)
    height = Inches(2)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table._graphic_frame.name = "test_table"

    # Set header
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(0, 2).text = "Status"

    return prs, slide


def test_plain_text_backward_compatibility():
    """Test that plain text works exactly as before (no regression)."""
    print("Testing backward compatibility with plain text...")

    prs, slide = create_test_presentation()

    # Find the text placeholder
    textbox = None
    for shape in slide.shapes:
        if shape.name == "test_placeholder":
            textbox = shape
            break

    assert textbox is not None, "Textbox not found"

    # Test 1: Simple plain text
    plain_text = "This is plain text with no markdown"
    result = populate_text_placeholder(textbox, plain_text)
    assert result == True, "populate_text_placeholder failed"
    assert textbox.text_frame.text == plain_text, "Text doesn't match"
    print("✓ Test 1 passed: Simple plain text")

    # Test 2: Plain text with newlines
    multiline_text = "Line 1\nLine 2\nLine 3"
    result = populate_text_placeholder(textbox, multiline_text)
    assert result == True, "populate_text_placeholder failed"
    assert textbox.text_frame.text == multiline_text, "Multiline text doesn't match"
    print("✓ Test 2 passed: Plain text with newlines")

    # Test 3: Plain text with special characters
    special_text = "Email: test@example.com | Price: $99.99"
    result = populate_text_placeholder(textbox, special_text)
    assert result == True, "populate_text_placeholder failed"
    assert textbox.text_frame.text == special_text, "Special characters text doesn't match"
    print("✓ Test 3 passed: Plain text with special characters")

    # Test 4: Empty string
    result = populate_text_placeholder(textbox, "")
    assert result == True, "populate_text_placeholder failed"
    assert textbox.text_frame.text == "", "Empty text doesn't match"
    print("✓ Test 4 passed: Empty string")

    # Test 5: Numbers and symbols
    numbers_text = "Numbers: 123 456 789 | Symbols: !@#$%^&*()"
    result = populate_text_placeholder(textbox, numbers_text)
    assert result == True, "populate_text_placeholder failed"
    assert textbox.text_frame.text == numbers_text, "Numbers/symbols text doesn't match"
    print("✓ Test 5 passed: Numbers and symbols")

    print("\n✅ All backward compatibility tests passed!")
    print("   Plain text functionality is unchanged.")


def test_table_backward_compatibility():
    """Test that plain text in tables works as before."""
    print("\nTesting backward compatibility with table population...")

    prs, slide = create_test_presentation()

    # Find the table
    table_shape = None
    for shape in slide.shapes:
        if shape.name == "test_table":
            table_shape = shape
            break

    assert table_shape is not None, "Table not found"

    # Plain text data
    plain_data = [
        ["Row 1", "Value 1", "Active"],
        ["Row 2", "Value 2", "Inactive"]
    ]

    result = populate_table(table_shape, plain_data, skip_header=True)
    assert result == True, "populate_table failed"

    # Verify data
    table = table_shape.table
    assert table.cell(1, 0).text == "Row 1", "Table cell 1,0 doesn't match"
    assert table.cell(1, 1).text == "Value 1", "Table cell 1,1 doesn't match"
    assert table.cell(2, 2).text == "Inactive", "Table cell 2,2 doesn't match"

    print("✓ Table population with plain text works correctly")
    print("\n✅ Table backward compatibility test passed!")


def test_markdown_formatting_works():
    """Test that markdown formatting actually works."""
    print("\nTesting markdown formatting (new functionality)...")

    prs, slide = create_test_presentation()

    # Find the text placeholder
    textbox = None
    for shape in slide.shapes:
        if shape.name == "test_placeholder":
            textbox = shape
            break

    # Test markdown text
    markdown_text = "**bold** and *italic* and __underline__"
    result = populate_text_placeholder(textbox, markdown_text)
    assert result == True, "populate_text_placeholder failed"

    # Verify formatting was applied (check number of runs)
    text_frame = textbox.text_frame
    para = text_frame.paragraphs[0]

    # Should have multiple runs (plain + bold + plain + italic + plain + underline)
    assert len(para.runs) > 1, "Markdown should create multiple runs"

    # Verify text content (without markdown markers)
    full_text = "".join(run.text for run in para.runs)
    expected_text = "bold and italic and underline"
    assert full_text == expected_text, f"Expected '{expected_text}', got '{full_text}'"

    # Verify first run is bold
    assert para.runs[0].font.bold == True, "First run should be bold"

    print("✓ Markdown formatting creates multiple runs")
    print("✓ Bold formatting is applied correctly")
    print("\n✅ Markdown formatting test passed!")


def test_performance():
    """Test that plain text performance is unchanged."""
    import time

    print("\nTesting performance (plain text should be fast)...")

    prs, slide = create_test_presentation()
    textbox = None
    for shape in slide.shapes:
        if shape.name == "test_placeholder":
            textbox = shape
            break

    # Test plain text performance
    plain_text = "Plain text without any markdown" * 100

    start = time.time()
    for i in range(100):
        populate_text_placeholder(textbox, plain_text)
    plain_elapsed = time.time() - start

    print(f"✓ 100 iterations of plain text: {plain_elapsed:.4f}s")
    print(f"✓ Average per iteration: {(plain_elapsed/100)*1000:.2f}ms")

    # Should be very fast (< 50ms per iteration on average)
    assert (plain_elapsed / 100) < 0.05, "Plain text performance regression detected"

    print("\n✅ Performance test passed!")
    print("   No performance regression for plain text.")


if __name__ == '__main__':
    try:
        test_plain_text_backward_compatibility()
        test_table_backward_compatibility()
        test_markdown_formatting_works()
        test_performance()

        print("\n" + "="*60)
        print("🎉 ALL TESTS PASSED!")
        print("="*60)
        print("\nSummary:")
        print("  ✅ Backward compatibility maintained")
        print("  ✅ Plain text works exactly as before")
        print("  ✅ Markdown formatting works correctly")
        print("  ✅ No performance regression")

    except AssertionError as e:
        print(f"\n❌ Test failed: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"\n❌ Unexpected error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
