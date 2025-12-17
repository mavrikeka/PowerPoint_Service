#!/usr/bin/env python3
"""
Extract data from PowerPoint presentations to JSON format.
This is the reverse operation of populate_ppt.py
"""

from pptx import Presentation
import json
import sys


def extract_text_from_shape(shape):
    """Extract text from a shape."""
    try:
        if hasattr(shape, "text_frame"):
            return shape.text_frame.text
        elif hasattr(shape, "text"):
            return shape.text
    except:
        pass
    return None


def extract_table_data(table_shape):
    """Extract data from a table shape."""
    if not table_shape.has_table:
        return None

    table = table_shape.table
    data = []

    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text)
        data.append(row_data)

    return data


def extract_data_from_presentation(pptx_path, slide_index=0, shape_names=None):
    """
    Extract data from a PowerPoint presentation.

    Args:
        pptx_path: Path to the PowerPoint file
        slide_index: Which slide to extract from (default: 0)
        shape_names: Optional list of shape names to extract. If None, extracts all.

    Returns:
        Dictionary with shape names as keys and extracted data as values
    """
    try:
        prs = Presentation(pptx_path)

        if len(prs.slides) == 0:
            return {"error": "No slides found in presentation"}

        if slide_index >= len(prs.slides):
            return {"error": f"Slide index {slide_index} out of range (presentation has {len(prs.slides)} slides)"}

        slide = prs.slides[slide_index]
        extracted_data = {}

        # Extract data from all shapes
        for shape in slide.shapes:
            shape_name = shape.name

            # Skip if we're filtering by shape names and this isn't in the list
            if shape_names and shape_name not in shape_names:
                continue

            # Extract table data
            if shape.has_table:
                extracted_data[shape_name] = extract_table_data(shape)
            else:
                # Extract text data
                text = extract_text_from_shape(shape)
                if text:
                    extracted_data[shape_name] = text

        return extracted_data

    except FileNotFoundError:
        return {"error": f"File not found: {pptx_path}"}
    except Exception as e:
        return {"error": f"Error extracting data: {str(e)}"}


def extract_all_slides(pptx_path):
    """Extract data from all slides in a presentation."""
    try:
        prs = Presentation(pptx_path)
        all_slides_data = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                "slide_index": slide_idx,
                "shapes": {}
            }

            for shape in slide.shapes:
                shape_name = shape.name

                if shape.has_table:
                    slide_data["shapes"][shape_name] = {
                        "type": "table",
                        "data": extract_table_data(shape)
                    }
                else:
                    text = extract_text_from_shape(shape)
                    if text:
                        slide_data["shapes"][shape_name] = {
                            "type": "text",
                            "data": text
                        }

            all_slides_data.append(slide_data)

        return all_slides_data

    except Exception as e:
        return {"error": f"Error extracting data: {str(e)}"}


def main():
    """Main function for command-line usage."""
    if len(sys.argv) < 2:
        print("Usage: python extract_from_ppt.py <path_to_pptx> [slide_index] [--all-slides]")
        print("\nExamples:")
        print("  python extract_from_ppt.py presentation.pptx")
        print("  python extract_from_ppt.py presentation.pptx 0")
        print("  python extract_from_ppt.py presentation.pptx --all-slides")
        sys.exit(1)

    pptx_path = sys.argv[1]

    # Check for --all-slides flag
    if "--all-slides" in sys.argv:
        print(f"Extracting data from all slides in: {pptx_path}\n")
        data = extract_all_slides(pptx_path)
    else:
        # Get slide index if provided
        slide_index = 0
        if len(sys.argv) > 2 and sys.argv[2].isdigit():
            slide_index = int(sys.argv[2])

        print(f"Extracting data from slide {slide_index} in: {pptx_path}\n")
        data = extract_data_from_presentation(pptx_path, slide_index)

    # Print as formatted JSON
    print(json.dumps(data, indent=2, ensure_ascii=False))

    # Optionally save to file
    output_file = pptx_path.replace('.pptx', '_extracted.json')
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

    print(f"\n✓ Data saved to: {output_file}")


if __name__ == "__main__":
    main()
